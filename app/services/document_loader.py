from pathlib import Path
from typing import List, Union
from langchain_core.documents import Document
from langchain_community.document_loaders.parsers import TesseractBlobParser
from langchain_community.document_loaders.blob_loaders.schema import Blob

import fitz  # PyMuPDF
import docx2txt
from pptx import Presentation
from PIL import Image

class UniversalDocumentLoader:
    def __init__(self, min_text_len=20):
        self.ocr_parser = TesseractBlobParser()
        self.min_text_len = min_text_len

    def load(self, input_path: Union[str, Path]) -> List[Document]:
        path = Path(input_path)
        files = [path] if path.is_file() else list(path.rglob("*.*"))
        all_docs = []

        for file in files:
            ext = file.suffix.lower()
            if ext == ".pdf":
                all_docs.extend(self._load_pdf(file))
            elif ext in [".jpg", ".jpeg", ".png"]:
                all_docs.extend(self._load_image(file))
            elif ext == ".docx":
                all_docs.extend(self._load_docx(file))
            elif ext == ".pptx":
                all_docs.extend(self._load_pptx(file))
            # You can extend this to .doc and .ppt via textract if needed

        return all_docs

    def _load_pdf(self, file: Path) -> List[Document]:
        doc = fitz.open(str(file))
        documents = []

        for i in range(len(doc)):
            page = doc.load_page(i)
            text = page.get_text().strip()

            if len(text) >= self.min_text_len:
                documents.append(Document(
                    page_content=text,
                    metadata={"source": str(file), "page": i + 1, "type": "text"}
                ))
            else:
                pix = page.get_pixmap(dpi=300)
                img_data = pix.tobytes("png")
                blob = Blob.from_data(img_data, path=f"{file}-page-{i+1}.png")
                ocr_docs = self.ocr_parser.parse(blob)

                for ocr_doc in ocr_docs:
                    ocr_doc.metadata.update({
                        "source": str(file),
                        "page": i + 1,
                        "type": "ocr"
                    })
                    documents.append(ocr_doc)

        return documents

    def _load_image(self, file: Path) -> List[Document]:
        with open(file, "rb") as f:
            blob = Blob.from_data(f.read(), path=str(file))
        return self.ocr_parser.parse(blob)

    def _load_docx(self, file: Path) -> List[Document]:
        text = docx2txt.process(str(file))
        return [Document(page_content=text.strip(), metadata={"source": str(file), "type": "docx"})] if text else []

    def _load_pptx(self, file: Path) -> List[Document]:
        prs = Presentation(str(file))
        documents = []

        for i, slide in enumerate(prs.slides):
            text = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
            full_text = "\n".join(text).strip()
            if full_text:
                documents.append(Document(
                    page_content=full_text,
                    metadata={"source": str(file), "slide": i + 1, "type": "pptx"}
                ))
        return documents
