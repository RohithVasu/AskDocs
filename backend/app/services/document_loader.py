from pathlib import Path
from typing import List, Union
import io
from langchain_core.documents import Document
from loguru import logger
import fitz  # PyMuPDF
import docx2txt
from pptx import Presentation
from PIL import Image
import pytesseract

class UniversalDocumentLoader:
    def __init__(self, min_text_len=20):
        self.min_text_len = min_text_len

    def _process_image(self, image_bytes: bytes) -> str:
        """Extract text from image using OCR."""
        try:
            image = Image.open(io.BytesIO(image_bytes))
            text = pytesseract.image_to_string(image)
            return text.strip()
        except Exception as e:
            logger.error(f"Error processing image with OCR: {e}")
            return ""

    def load(self, input_path: Union[str, Path]) -> List[Document]:
        path = Path(input_path)
        files = [path] if path.is_file() else list(path.rglob("*.*"))
        all_docs = []

        for file in files:
            try:
                ext = file.suffix.lower()
                if ext == ".pdf":
                    all_docs.extend(self._load_pdf(file))
                elif ext in [".jpg", ".jpeg", ".png"]:
                    all_docs.extend(self._load_image(file))
                elif ext == ".docx":
                    all_docs.extend(self._load_docx(file))
                elif ext == ".pptx":
                    all_docs.extend(self._load_pptx(file))
            except Exception as e:
                logger.error(f"Error loading file {file}: {e}")

        return all_docs

    def _load_pdf(self, file: Path) -> List[Document]:
        doc = fitz.open(str(file))
        documents = []

        for i in range(len(doc)):
            page = doc.load_page(i)
            text = page.get_text().strip()

            # If text is sufficient, use it
            if len(text) >= self.min_text_len:
                documents.append(Document(
                    page_content=text,
                    metadata={
                        "source": str(file.name),
                        "page_number": i + 1,
                        "content_category": "text",
                        "file_type": "pdf"
                    }
                ))
            
            # Check for images on the page
            image_list = page.get_images(full=True)
            if image_list:
                logger.info(f"Found {len(image_list)} images on page {i+1} of {file}")
                for img_index, img in enumerate(image_list):
                    try:
                        xref = img[0]
                        base_image = doc.extract_image(xref)
                        image_bytes = base_image["image"]
                        
                        # Extract text from image with OCR
                        ocr_text = self._process_image(image_bytes)
                        
                        if ocr_text:
                            documents.append(Document(
                                page_content=f"[Image OCR]: {ocr_text}",
                                metadata={
                                    "source": str(file.name),
                                    "page_number": i + 1,
                                    "content_category": "image",
                                    "file_type": "pdf",
                                    "image_index": img_index
                                }
                            ))
                    except Exception as e:
                        logger.error(f"Error extracting image {img_index} from PDF page {i+1}: {e}")

        return documents

    def _load_image(self, file: Path) -> List[Document]:
        with open(file, "rb") as f:
            image_bytes = f.read()
            
        ocr_text = self._process_image(image_bytes)
        
        if ocr_text:
            return [Document(
                page_content=ocr_text,
                metadata={
                    "source": str(file.name),
                    "content_category": "image",
                    "file_type": "image",
                    "page_number": 1
                }
            )]
        return []

    def _load_docx(self, file: Path) -> List[Document]:
        # Load text
        text = docx2txt.process(str(file))
        documents = [Document(
            page_content=text.strip(),
            metadata={
                "source": str(file.name),
                "content_category": "text",
                "file_type": "docx",
                "page_number": 1
            }
        )] if text else []
        
        # Load images using python-docx
        try:
            from docx import Document as DocxDocument
            doc = DocxDocument(str(file))
            
            for i, rel in enumerate(doc.part.rels.values()):
                if "image" in rel.target_ref:
                    image_part = rel.target_part
                    image_bytes = image_part.blob
                    
                    ocr_text = self._process_image(image_bytes)
                    if ocr_text:
                        documents.append(Document(
                            page_content=f"[Image OCR]: {ocr_text}",
                            metadata={
                                "source": str(file.name),
                                "content_category": "image",
                                "file_type": "docx",
                                "page_number": 1, # DOCX doesn't have clear pages
                                "image_index": i
                            }
                        ))
        except Exception as e:
            logger.error(f"Error extracting images from DOCX {file}: {e}")
            
        return documents

    def _load_pptx(self, file: Path) -> List[Document]:
        prs = Presentation(str(file))
        documents = []

        for i, slide in enumerate(prs.slides):
            text = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
                
                # Check for images
                if hasattr(shape, "image"):
                    try:
                        image_bytes = shape.image.blob
                        ocr_text = self._process_image(image_bytes)
                        if ocr_text:
                            documents.append(Document(
                                page_content=f"[Image OCR]: {ocr_text}",
                                metadata={
                                    "source": str(file.name),
                                    "page_number": i + 1,
                                    "content_category": "image",
                                    "file_type": "pptx"
                                }
                            ))
                    except Exception as e:
                        logger.error(f"Error processing image in PPTX {file} slide {i+1}: {e}")

            full_text = "\n".join(text).strip()
            if full_text:
                documents.append(Document(
                    page_content=full_text,
                    metadata={
                        "source": str(file.name),
                        "page_number": i + 1,
                        "content_category": "text",
                        "file_type": "pptx"
                    }
                ))
        return documents
